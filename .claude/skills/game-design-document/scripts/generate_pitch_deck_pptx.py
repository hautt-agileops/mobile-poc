#!/usr/bin/env python3
"""Render game_config.json's pitch section into a .pptx pitch deck.

Reads config.pitch.slides = [{title, bullets[], notes}]. Falls back to a
title slide + one slide per top-level GDD section heading if pitch is absent.

Usage:
    python scripts/generate_pitch_deck_pptx.py --config game_config.json --output "Title_Pitch_v01.pptx"

Needs: python-pptx==1.0.2
"""
import argparse

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

from _config import load_config, meta_line

ACCENT = RGBColor(0x1F, 0x4E, 0x79)
MUTED = RGBColor(0x60, 0x60, 0x60)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)


def _title_slide(prs, meta):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = ACCENT
    box = slide.shapes.add_textbox(Inches(0.7), Inches(2.2),
                                   Inches(8.6), Inches(3))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = meta["title"]
    r.font.size = Pt(44)
    r.font.bold = True
    r.font.color.rgb = WHITE
    sub = meta.get("subtitle") or meta.get("tagline")
    if sub:
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run()
        r2.text = sub
        r2.font.size = Pt(20)
        r2.font.italic = True
        r2.font.color.rgb = WHITE
    p3 = tf.add_paragraph()
    p3.alignment = PP_ALIGN.CENTER
    r3 = p3.add_run()
    r3.text = meta_line(meta)
    r3.font.size = Pt(14)
    r3.font.color.rgb = WHITE


def _content_slide(prs, title, bullets, notes=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tbox = slide.shapes.add_textbox(Inches(0.6), Inches(0.4),
                                    Inches(8.8), Inches(1))
    tp = tbox.text_frame.paragraphs[0]
    tr = tp.add_run()
    tr.text = title
    tr.font.size = Pt(30)
    tr.font.bold = True
    tr.font.color.rgb = ACCENT

    if bullets:
        bbox = slide.shapes.add_textbox(Inches(0.7), Inches(1.6),
                                        Inches(8.6), Inches(5))
        bf = bbox.text_frame
        bf.word_wrap = True
        for i, b in enumerate(bullets):
            para = bf.paragraphs[0] if i == 0 else bf.add_paragraph()
            run = para.add_run()
            run.text = f"•  {b}"
            run.font.size = Pt(18)
            para.space_after = Pt(8)
    if notes:
        slide.notes_slide.notes_text_frame.text = str(notes)


def build(cfg, output):
    prs = Presentation()  # 4:3 default; fine for pitch
    meta = cfg["meta"]
    _title_slide(prs, meta)

    slides = (cfg.get("pitch") or {}).get("slides")
    if slides:
        for s in slides:
            _content_slide(prs, s.get("title", ""),
                           s.get("bullets", []), s.get("notes"))
    else:
        # fallback: one slide per GDD section, first bullet = its body head
        for sec in cfg["sections"]:
            body = str(sec.get("body", "")).strip().split("\n\n")
            bullets = sec.get("bullets") or ([body[0]] if body and body[0] else [])
            _content_slide(prs, sec.get("heading", "Untitled"), bullets)

    prs.save(output)
    print(f"wrote {output}")


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--config", required=True)
    ap.add_argument("--output", required=True)
    args = ap.parse_args()
    build(load_config(args.config), args.output)


if __name__ == "__main__":
    main()
